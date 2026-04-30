from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt


OUT = Path(r"D:\GenAI_Project_2025\myHealthCoach\docs\aiagent_similar.pptx")

BRAND_BG = RGBColor(246, 248, 252)
BRAND_ACCENT = RGBColor(12, 61, 130)
BRAND_TEXT = RGBColor(34, 34, 34)
BRAND_MUTED = RGBColor(90, 90, 90)

FONT_TITLE = "Segoe UI Semibold"
FONT_BODY = "Calibri"

PRESENTER_NAME = "Code Insights Team"
PRESENTER_ROLE = "AI and Agentic AI Instructor"
PRESENTER_EMAIL = "codeinsights.ai@gmail.com"


def _apply_brand_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BRAND_BG

    banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.45))
    banner.fill.solid()
    banner.fill.fore_color.rgb = BRAND_ACCENT
    banner.line.fill.background()


def _style_title(shape) -> None:
    tf = shape.text_frame
    for p in tf.paragraphs:
        for run in p.runs:
            run.font.name = FONT_TITLE
            run.font.color.rgb = BRAND_ACCENT
            run.font.size = Pt(34)


def _add_footer(slide, page_no: int, total_pages: int) -> None:
    footer = slide.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.35))
    tf = footer.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = f"{PRESENTER_NAME} | {PRESENTER_ROLE} | {PRESENTER_EMAIL}"
    p.font.name = FONT_BODY
    p.font.size = Pt(12)
    p.font.color.rgb = BRAND_MUTED

    page_box = slide.shapes.add_textbox(Inches(12.3), Inches(7.08), Inches(0.8), Inches(0.35))
    p_tf = page_box.text_frame
    p_tf.clear()
    pp = p_tf.paragraphs[0]
    pp.text = f"{page_no}/{total_pages}"
    pp.font.name = FONT_BODY
    pp.font.size = Pt(12)
    pp.font.color.rgb = BRAND_MUTED


def _add_diagram_placeholder(slide, label: str) -> None:
    ph = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.2), Inches(1.4), Inches(5.6), Inches(4.8))
    ph.fill.solid()
    ph.fill.fore_color.rgb = RGBColor(255, 255, 255)
    ph.line.color.rgb = BRAND_ACCENT
    ph.line.width = Pt(1.5)

    tf = ph.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = label
    p.font.name = FONT_BODY
    p.font.bold = True
    p.font.size = Pt(18)
    p.font.color.rgb = BRAND_ACCENT


def add_title(prs: Presentation, title: str, subtitle: str, page_no: int, total_pages: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    _apply_brand_background(slide)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    _style_title(slide.shapes.title)

    sub_tf = slide.placeholders[1].text_frame
    for p in sub_tf.paragraphs:
        for run in p.runs:
            run.font.name = FONT_BODY
            run.font.size = Pt(20)
            run.font.color.rgb = BRAND_TEXT

    _add_footer(slide, page_no, total_pages)


def add_bullets(
    prs: Presentation,
    title: str,
    bullets: list[str],
    page_no: int,
    total_pages: int,
    diagram_label: str | None = None,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _apply_brand_background(slide)
    slide.shapes.title.text = title

    _style_title(slide.shapes.title)

    text_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.35), Inches(6.2), Inches(5.7))
    tf = text_box.text_frame
    tf.clear()
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.name = FONT_BODY
        p.font.size = Pt(22)
        p.font.color.rgb = BRAND_TEXT

    if diagram_label:
        _add_diagram_placeholder(slide, diagram_label)

    _add_footer(slide, page_no, total_pages)


def build() -> Presentation:
    prs = Presentation()

    total_pages = 12

    add_title(
        prs,
        "Agentic AI",
        "Foundations, Frameworks, and Practical Classroom View",
        page_no=1,
        total_pages=total_pages,
    )

    add_bullets(
        prs,
        "About Me",
        [
            "A passionate technocrat",
            "AI research and applied innovation",
            "Automation and intelligent systems",
            "Author, educator, and mentor",
            "Adventurer and continuous learner",
        ],
        page_no=2,
        total_pages=total_pages,
    )

    add_bullets(
        prs,
        "Agenda",
        [
            "Definitions: AI Agent and Agentic AI",
            "Types of agents and when to use each",
            "Frameworks: LangChain, LangGraph, LangSmith",
            "Patterns: ReAct, Reflection, Plan-and-Execute",
            "Case study: MyHealthCoach implementation",
        ],
        page_no=3,
        total_pages=total_pages,
        diagram_label="Diagram Placeholder\n(Agenda Flow Graphic)",
    )

    add_bullets(
        prs,
        "Generative AI vs AI Agents vs Agentic AI",
        [
            "Generative AI creates outputs (text, code, media)",
            "AI Agents perform actions with tools and APIs",
            "Agentic AI plans, coordinates, and adapts execution",
            "Modern systems combine all three for outcomes",
        ],
        page_no=4,
        total_pages=total_pages,
        diagram_label="Image Placeholder\n(Concept Comparison Diagram)",
    )

    add_bullets(
        prs,
        "Relationship: Generative AI, Agentic AI, AI Agents",
        [
            "Agentic AI provides the brain and the plan",
            "AI Agents are the hands and feet that execute",
            "Generative AI provides raw creative output",
            "A mature system combines all three",
        ],
        page_no=5,
        total_pages=total_pages,
        diagram_label="Diagram Placeholder\n(Relationship Triangle)",
    )

    add_bullets(
        prs,
        "Table: Framework and Pattern Mapping",
        [
            "LangChain: model, prompt, tools composition",
            "LangGraph: stateful orchestration and control",
            "LangSmith: tracing, experiments, evaluations",
            "ReAct and Reflection: reasoning behavior patterns",
        ],
        page_no=6,
        total_pages=total_pages,
        diagram_label="Table/Image Placeholder\n(Framework Matrix)",
    )

    add_bullets(
        prs,
        "ReAct Framework",
        [
            "Reason -> Act -> Observe iterative loop",
            "Best for tool-use and uncertain contexts",
            "Improves grounding and task completion quality",
            "Tradeoff: increased latency and token consumption",
        ],
        page_no=7,
        total_pages=total_pages,
        diagram_label="Diagram Placeholder\n(ReAct Loop)",
    )

    add_bullets(
        prs,
        "Reflection Framework",
        [
            "Generate draft, critique, then revise response",
            "Detect gaps in logic, evidence, and safety",
            "Useful for high-stakes or nuanced domains",
            "Tradeoff: extra step adds compute time",
        ],
        page_no=8,
        total_pages=total_pages,
        diagram_label="Diagram Placeholder\n(Reflection Critique Cycle)",
    )

    add_bullets(
        prs,
        "LangChain, LangGraph, LangSmith",
        [
            "LangChain: tool and prompt building blocks",
            "LangGraph: deterministic orchestration with state",
            "LangSmith: trace-level visibility and evaluations",
            "Stack enables production-ready agentic systems",
        ],
        page_no=9,
        total_pages=total_pages,
        diagram_label="Architecture Placeholder\n(Framework Stack)",
    )

    add_bullets(
        prs,
        "MyHealthCoach Solution Overview",
        [
            "Streamlit UI captures profile + health questions",
            "Supervisor ReAct routes to specialist advisors",
            "Doctor lookup: NPI (US), Nominatim/Overpass (India)",
            "Azure-ready architecture with traceable execution",
        ],
        page_no=10,
        total_pages=total_pages,
        diagram_label="Diagram Placeholder\n(MyHealthCoach Architecture)",
    )

    add_bullets(
        prs,
        "Implementation and Evaluation",
        [
            "Live scoring with weighted criteria",
            "Batch regression with 20 fixed prompts",
            "LangSmith feedback mapping and run_id submission",
            "Continuous quality loop for classroom demos and production",
        ],
        page_no=11,
        total_pages=total_pages,
        diagram_label="Image Placeholder\n(Evaluation Pipeline)",
    )

    add_bullets(
        prs,
        "Thank You",
        [
            "Questions and discussion",
            "Contact: update presenter details in the generator script",
            "Use this deck as your classroom base and customize examples",
        ],
        page_no=12,
        total_pages=total_pages,
    )

    return prs


if __name__ == "__main__":
    prs = build()
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Created: {OUT}")
