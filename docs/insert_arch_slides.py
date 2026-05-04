"""
Inserts two new architecture diagram slides into MyHealthCoach_YouTube_Video.pptx:

  After Slide  6  (Supervisor + Specialist Architecture)
      → NEW: "Data Flow Architecture" with dfd_mhc.drawio.png

  After Slide 16  (Azure Service Integrations section)  [+1 offset after first insert]
      → NEW: "Technical Architecture" with tech_mhc.drawio.png

Each new slide has:
  - Dark gradient background matching the deck palette
  - Slide title (white, top)
  - Diagram image (centred, fills most of slide)
  - logo.png watermark bottom-right

Run: python docs/insert_arch_slides.py
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import copy

ROOT   = Path(__file__).resolve().parent.parent
PPTX   = ROOT / "docs" / "MyHealthCoach_YouTube_Video.pptx"
DFD    = ROOT / "docs" / "images" / "dfd_mhc.drawio.png"
TECH   = ROOT / "docs" / "images" / "tech_mhc.drawio.png"
LOGO   = ROOT / "assets" / "logo.png"

# Deck colour palette
DARK_BLUE   = RGBColor(0x0C, 0x3D, 0x82)
MID_BLUE    = RGBColor(0x00, 0x70, 0xC0)
LIGHT_BLUE  = RGBColor(0xEB, 0xF2, 0xFF)
ACCENT_GREEN= RGBColor(0x00, 0x99, 0x4C)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation(str(PPTX))
SW  = prs.slide_width    # 9144000 EMU  (13.33 in for widescreen 16:9)
SH  = prs.slide_height   # 5143500 EMU  (7.5 in)

# ── pick the "blank" layout (index 6 is typically blank in Office themes) ─────
# Fall back to the layout with fewest placeholders
def blank_layout(prs):
    best = None
    for layout in prs.slide_layouts:
        ph_count = len(layout.placeholders)
        if best is None or ph_count < len(best.placeholders):
            best = layout
    return best

layout = blank_layout(prs)


# ── helper: move slide from end to a given position ───────────────────────────
def move_slide_to(prs, slide, position):
    """Move the most-recently-added slide to `position` (0-indexed)."""
    xml_slides = prs.slides._sldIdLst
    last = xml_slides[-1]
    xml_slides.remove(last)
    xml_slides.insert(position, last)


# ── helper: solid background fill on slide ────────────────────────────────────
def set_slide_bg(slide, rgb: RGBColor):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = rgb


# ── helper: add a text box ────────────────────────────────────────────────────
def add_textbox(slide, text, left, top, width, height,
                font_size=28, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


# ── helper: build one architecture slide ─────────────────────────────────────
def build_arch_slide(title_text: str, diagram_path: Path, insert_at: int):
    slide = prs.slides.add_slide(layout)
    move_slide_to(prs, slide, insert_at)

    # Background — dark navy
    set_slide_bg(slide, DARK_BLUE)

    # Thin accent bar at top
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(0), Inches(0),
        SW, Inches(0.07),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_GREEN
    bar.line.fill.background()

    # Title
    add_textbox(
        slide, title_text,
        left=Inches(0.35), top=Inches(0.1),
        width=SW - Inches(3), height=Inches(0.7),
        font_size=30, bold=True, color=WHITE,
    )

    # Subtitle / description tag (right-aligned)
    tag = "Architecture Diagram"
    add_textbox(
        slide, tag,
        left=SW - Inches(2.8), top=Inches(0.12),
        width=Inches(2.6), height=Inches(0.55),
        font_size=13, bold=False,
        color=RGBColor(0xAA, 0xCC, 0xFF),
        align=PP_ALIGN.RIGHT,
    )

    # Diagram image — centred, leaving margin for title and logo
    IMG_MARGIN_L = Inches(0.25)
    IMG_TOP      = Inches(0.82)
    IMG_W        = SW - Inches(0.5)
    IMG_H        = SH - Inches(1.35)   # leave room for logo at bottom
    slide.shapes.add_picture(
        str(diagram_path),
        IMG_MARGIN_L, IMG_TOP,
        IMG_W, IMG_H,
    )

    # Logo watermark — bottom-right
    LOGO_W = Inches(1.8)
    LOGO_H = Inches(0.5)
    LOGO_M = Inches(0.15)
    slide.shapes.add_picture(
        str(LOGO),
        SW - LOGO_W - LOGO_M,
        SH - LOGO_H - LOGO_M,
        LOGO_W, LOGO_H,
    )

    print(f"✔  Inserted '{title_text}' at position {insert_at + 1}")


# ── Insert slide 1: Data Flow Architecture after slide 6 ─────────────────────
build_arch_slide(
    title_text   = "Data Flow Architecture — Request to Response",
    diagram_path = DFD,
    insert_at    = 6,   # 0-indexed → after original slide 6 (index 5)
)

# ── Insert slide 2: Technical Architecture after original slide 16 ────────────
# Original slide 16 is now at index 17 (shifted +1 by first insert)
build_arch_slide(
    title_text   = "Technical Architecture — Azure & AI Stack",
    diagram_path = TECH,
    insert_at    = 17,
)

# ── Save ──────────────────────────────────────────────────────────────────────
prs.save(str(PPTX))
print(f"\nSaved: {PPTX}")
print(f"Total slides now: {len(prs.slides)}")
