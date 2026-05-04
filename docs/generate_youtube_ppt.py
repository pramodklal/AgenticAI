from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "MyHealthCoach_YouTube_Video.pptx"

# ── brand palette ──────────────────────────────────────────────────────────────
DARK_BLUE = RGBColor(12, 61, 130)
MID_BLUE = RGBColor(0, 112, 192)
LIGHT_BLUE = RGBColor(235, 242, 255)
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(34, 34, 34)
MID_GRAY = RGBColor(100, 100, 100)
ACCENT_GREEN = RGBColor(0, 153, 76)
ACCENT_AMBER = RGBColor(255, 153, 0)


# ── helpers ────────────────────────────────────────────────────────────────────

def _set_font(run, size_pt: int, bold: bool = False, color: RGBColor = DARK_GRAY):
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color


def add_title_slide(prs: Presentation, title: str, subtitle: str, tag: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # full-bleed background
    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.fill.background()

    # accent stripe
    stripe = slide.shapes.add_shape(1, Inches(0), Inches(5.4), Inches(13.33), Inches(0.18))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = ACCENT_AMBER
    stripe.line.fill.background()

    # title
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.7), Inches(1.8))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    _set_font(run, 40, bold=True, color=WHITE)

    # subtitle
    tb2 = slide.shapes.add_textbox(Inches(0.8), Inches(3.3), Inches(11.7), Inches(1.4))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = subtitle
    _set_font(run2, 22, color=RGBColor(200, 220, 255))

    # channel tag
    tb3 = slide.shapes.add_textbox(Inches(0.8), Inches(5.8), Inches(11.7), Inches(0.6))
    tf3 = tb3.text_frame
    p3 = tf3.paragraphs[0]
    run3 = p3.add_run()
    run3.text = tag
    _set_font(run3, 14, color=ACCENT_AMBER)


def add_section_divider(prs: Presentation, section_number: str, section_title: str, subtitle: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = MID_BLUE
    bg.line.fill.background()

    stripe = slide.shapes.add_shape(1, Inches(0), Inches(3.6), Inches(13.33), Inches(0.12))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = WHITE
    stripe.line.fill.background()

    tb_num = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.7), Inches(0.8))
    p_num = tb_num.text_frame.paragraphs[0]
    r_num = p_num.add_run()
    r_num.text = section_number
    _set_font(r_num, 16, color=ACCENT_AMBER)

    tb = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.7), Inches(1.4))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = section_title
    _set_font(run, 36, bold=True, color=WHITE)

    if subtitle:
        tb2 = slide.shapes.add_textbox(Inches(0.8), Inches(3.8), Inches(11.7), Inches(0.9))
        p2 = tb2.text_frame.paragraphs[0]
        r2 = p2.add_run()
        r2.text = subtitle
        _set_font(r2, 20, color=RGBColor(210, 228, 255))


def add_bullet_slide(
    prs: Presentation,
    title: str,
    bullets: list[str],
    footer: str = "",
    sub_title: str = "",
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # header bar
    hdr = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.15))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = DARK_BLUE
    hdr.line.fill.background()

    tb_title = slide.shapes.add_textbox(Inches(0.5), Inches(0.18), Inches(12.3), Inches(0.8))
    tf_title = tb_title.text_frame
    p_title = tf_title.paragraphs[0]
    r_title = p_title.add_run()
    r_title.text = title
    _set_font(r_title, 26, bold=True, color=WHITE)

    if sub_title:
        tb_sub = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.3), Inches(0.5))
        p_sub = tb_sub.text_frame.paragraphs[0]
        r_sub = p_sub.add_run()
        r_sub.text = sub_title
        _set_font(r_sub, 15, color=MID_BLUE)

    content_top = Inches(1.85) if sub_title else Inches(1.35)

    tb_body = slide.shapes.add_textbox(Inches(0.5), content_top, Inches(12.3), Inches(4.6))
    tf_body = tb_body.text_frame
    tf_body.word_wrap = True

    for i, bullet in enumerate(bullets):
        p = tf_body.paragraphs[0] if i == 0 else tf_body.add_paragraph()
        r = p.add_run()
        if bullet.startswith(">>"):
            r.text = "    " + bullet[2:].strip()
            _set_font(r, 17, color=MID_GRAY)
        else:
            r.text = "▸  " + bullet
            _set_font(r, 19, color=DARK_GRAY)
        p.space_after = Pt(6)

    if footer:
        ftr = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.45))
        pf = ftr.text_frame.paragraphs[0]
        rf = pf.add_run()
        rf.text = footer
        _set_font(rf, 12, color=MID_GRAY)


def add_two_col_slide(
    prs: Presentation,
    title: str,
    left_heading: str,
    left_bullets: list[str],
    right_heading: str,
    right_bullets: list[str],
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    hdr = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.15))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = DARK_BLUE
    hdr.line.fill.background()

    tb_title = slide.shapes.add_textbox(Inches(0.5), Inches(0.18), Inches(12.3), Inches(0.8))
    r_title = tb_title.text_frame.paragraphs[0].add_run()
    r_title.text = title
    _set_font(r_title, 26, bold=True, color=WHITE)

    # left column
    left_bg = slide.shapes.add_shape(1, Inches(0.35), Inches(1.25), Inches(5.9), Inches(5.8))
    left_bg.fill.solid()
    left_bg.fill.fore_color.rgb = LIGHT_BLUE
    left_bg.line.color.rgb = DARK_BLUE

    tb_lh = slide.shapes.add_textbox(Inches(0.55), Inches(1.35), Inches(5.5), Inches(0.5))
    r_lh = tb_lh.text_frame.paragraphs[0].add_run()
    r_lh.text = left_heading
    _set_font(r_lh, 16, bold=True, color=DARK_BLUE)

    tb_lb = slide.shapes.add_textbox(Inches(0.55), Inches(1.9), Inches(5.5), Inches(4.8))
    tb_lb.text_frame.word_wrap = True
    for i, b in enumerate(left_bullets):
        p = tb_lb.text_frame.paragraphs[0] if i == 0 else tb_lb.text_frame.add_paragraph()
        r = p.add_run()
        r.text = "▸  " + b
        _set_font(r, 17, color=DARK_GRAY)
        p.space_after = Pt(5)

    # right column
    right_bg = slide.shapes.add_shape(1, Inches(7.05), Inches(1.25), Inches(5.9), Inches(5.8))
    right_bg.fill.solid()
    right_bg.fill.fore_color.rgb = RGBColor(255, 250, 235)
    right_bg.line.color.rgb = ACCENT_AMBER

    tb_rh = slide.shapes.add_textbox(Inches(7.25), Inches(1.35), Inches(5.5), Inches(0.5))
    r_rh = tb_rh.text_frame.paragraphs[0].add_run()
    r_rh.text = right_heading
    _set_font(r_rh, 16, bold=True, color=RGBColor(160, 90, 0))

    tb_rb = slide.shapes.add_textbox(Inches(7.25), Inches(1.9), Inches(5.5), Inches(4.8))
    tb_rb.text_frame.word_wrap = True
    for i, b in enumerate(right_bullets):
        p = tb_rb.text_frame.paragraphs[0] if i == 0 else tb_rb.text_frame.add_paragraph()
        r = p.add_run()
        r.text = "▸  " + b
        _set_font(r, 17, color=DARK_GRAY)
        p.space_after = Pt(5)


def add_code_slide(prs: Presentation, title: str, code_lines: list[str], note: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    hdr = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.15))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = DARK_BLUE
    hdr.line.fill.background()

    tb_title = slide.shapes.add_textbox(Inches(0.5), Inches(0.18), Inches(12.3), Inches(0.8))
    r_title = tb_title.text_frame.paragraphs[0].add_run()
    r_title.text = title
    _set_font(r_title, 26, bold=True, color=WHITE)

    code_bg = slide.shapes.add_shape(1, Inches(0.4), Inches(1.3), Inches(12.5), Inches(5.2))
    code_bg.fill.solid()
    code_bg.fill.fore_color.rgb = RGBColor(30, 30, 40)
    code_bg.line.color.rgb = RGBColor(80, 80, 120)

    tb_code = slide.shapes.add_textbox(Inches(0.65), Inches(1.45), Inches(12.1), Inches(4.9))
    tf_code = tb_code.text_frame
    tf_code.word_wrap = False
    for i, line in enumerate(code_lines):
        p = tf_code.paragraphs[0] if i == 0 else tf_code.add_paragraph()
        r = p.add_run()
        r.text = line
        r.font.size = Pt(15)
        r.font.bold = False
        r.font.color.rgb = RGBColor(180, 220, 180)
        r.font.name = "Courier New"

    if note:
        tb_note = slide.shapes.add_textbox(Inches(0.5), Inches(6.65), Inches(12.3), Inches(0.55))
        pn = tb_note.text_frame.paragraphs[0]
        rn = pn.add_run()
        rn.text = "💡  " + note
        _set_font(rn, 13, color=ACCENT_AMBER)


def add_chat_demo_slide(
    prs: Presentation,
    title: str,
    user_text: str,
    agent_label: str,
    agent_lines: list[str],
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    hdr = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(1.15))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = DARK_BLUE
    hdr.line.fill.background()

    tb_title = slide.shapes.add_textbox(Inches(0.5), Inches(0.18), Inches(12.3), Inches(0.8))
    r_title = tb_title.text_frame.paragraphs[0].add_run()
    r_title.text = title
    _set_font(r_title, 26, bold=True, color=WHITE)

    # user bubble (right-aligned)
    ub = slide.shapes.add_shape(4, Inches(5.5), Inches(1.3), Inches(7.3), Inches(0.85))
    ub.fill.solid()
    ub.fill.fore_color.rgb = DARK_BLUE
    ub.line.fill.background()
    tbu = ub.text_frame
    tbu.word_wrap = True
    pu = tbu.paragraphs[0]
    ru = pu.add_run()
    ru.text = "User:  " + user_text
    _set_font(ru, 16, bold=True, color=WHITE)

    # routing indicator
    tb_route = slide.shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(12.3), Inches(0.4))
    pr = tb_route.text_frame.paragraphs[0]
    rr = pr.add_run()
    rr.text = "Supervisor → routes to → " + agent_label
    _set_font(rr, 14, color=MID_BLUE)

    # agent bubble
    ab = slide.shapes.add_shape(4, Inches(0.5), Inches(2.85), Inches(12.4), Inches(4.1))
    ab.fill.solid()
    ab.fill.fore_color.rgb = LIGHT_BLUE
    ab.line.color.rgb = MID_BLUE
    tba = ab.text_frame
    tba.word_wrap = True

    p0 = tba.paragraphs[0]
    r0 = p0.add_run()
    r0.text = "MyHealthCoach:"
    _set_font(r0, 15, bold=True, color=DARK_BLUE)

    for line in agent_lines:
        p = tba.add_paragraph()
        r = p.add_run()
        r.text = line
        _set_font(r, 16, color=DARK_GRAY)
        p.space_before = Pt(4)


def add_closing_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.fill.background()

    stripe = slide.shapes.add_shape(1, Inches(0), Inches(5.0), Inches(13.33), Inches(0.18))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = ACCENT_AMBER
    stripe.line.fill.background()

    tb1 = slide.shapes.add_textbox(Inches(1.0), Inches(1.0), Inches(11.3), Inches(1.4))
    r1 = tb1.text_frame.paragraphs[0].add_run()
    r1.text = "Thanks for watching!"
    _set_font(r1, 38, bold=True, color=WHITE)

    tb2 = slide.shapes.add_textbox(Inches(1.0), Inches(2.6), Inches(11.3), Inches(0.7))
    r2 = tb2.text_frame.paragraphs[0].add_run()
    r2.text = "Like · Subscribe · Comment below with your questions"
    _set_font(r2, 20, color=RGBColor(200, 220, 255))

    tb3 = slide.shapes.add_textbox(Inches(1.0), Inches(3.6), Inches(11.3), Inches(0.6))
    r3 = tb3.text_frame.paragraphs[0].add_run()
    r3.text = "GitHub:  github.com/pramodklal/AgenticAI"
    _set_font(r3, 18, bold=True, color=ACCENT_AMBER)

    links = [
        "Source code:  Full Python project — LangGraph + Azure + Streamlit",
        "pip install -r requirements.txt  |  streamlit run app.py",
    ]
    for idx, line in enumerate(links):
        tb = slide.shapes.add_textbox(Inches(1.0), Inches(4.35 + idx * 0.55), Inches(11.3), Inches(0.5))
        r = tb.text_frame.paragraphs[0].add_run()
        r.text = line
        _set_font(r, 15, color=RGBColor(210, 228, 255))


# ── deck builder ───────────────────────────────────────────────────────────────

def build_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # ── SLIDE 1: Title ─────────────────────────────────────────────────────────
    add_title_slide(
        prs,
        "Build an AI Health Coach with Agentic AI",
        "LangGraph · Azure OpenAI · Streamlit · LangSmith\nSupervisor + Specialist Multi-Agent Architecture in Python",
        "github.com/pramodklal/AgenticAI",
    )

    # ── SLIDE 2: Agenda ────────────────────────────────────────────────────────
    add_bullet_slide(
        prs,
        "What We Cover in This Video",
        [
            "What is Agentic AI — and why it matters",
            "ReAct pattern and Supervisor + Specialist architecture",
            "LangGraph state machine and graph workflow code",
            "Domain-scoped specialist agents with tools",
            "Streamlit UI, user profile, and session state",
            "Azure OpenAI · AI Search · Cosmos DB integrations",
            "LangSmith tracing and custom evaluation scoring",
            "Doctor lookup — US (NPI Registry) and India (Overpass API)",
            "Batch regression testing and feedback pipeline",
            "Live demo: pre-diabetic patient query end-to-end",
        ],
        footer="Duration: ~22 minutes  |  Code: github.com/pramodklal/AgenticAI",
    )

    # ── SLIDE 3: Section — Agentic AI Concepts ─────────────────────────────────
    add_section_divider(prs, "SECTION 01", "Agentic AI Concepts", "From single LLM call to multi-step orchestration")

    # ── SLIDE 4: Regular AI vs Agentic AI ─────────────────────────────────────
    add_two_col_slide(
        prs,
        "Regular AI Call vs Agentic AI",
        "Regular LLM Call",
        [
            "User sends a message",
            "Model returns a single text answer",
            "One step — no tools, no state",
            "Limited by training data only",
            "No observability of reasoning",
        ],
        "Agentic AI System",
        [
            "User sends a message",
            "Supervisor reasons about intent",
            "Delegates to the right specialist",
            "Specialist calls tools for real data",
            "Synthesizes grounded, safe response",
            "Every step is traced and scored",
        ],
    )

    # ── SLIDE 5: ReAct Pattern ─────────────────────────────────────────────────
    add_bullet_slide(
        prs,
        "ReAct Pattern — Reason, Act, Observe",
        [
            "ReAct = Reasoning + Acting in an interleaved loop",
            "Step 1 — Reason:  What information do I need to answer this?",
            "Step 2 — Act:  Call the appropriate tool (sleep summary, nutrition docs…)",
            "Step 3 — Observe:  Integrate the tool result into reasoning context",
            "Repeat until enough evidence is gathered",
            "Final step — Synthesize:  Generate structured, grounded response",
            "Benefit: explainable, evidence-backed advice instead of hallucinated text",
        ],
        footer="Used in both supervisor and all 6 specialist agents in MyHealthCoach",
    )

    # ── SLIDE 6: Supervisor + Specialist Pattern ───────────────────────────────
    add_bullet_slide(
        prs,
        "Supervisor + Specialist Architecture",
        [
            "One supervisor agent handles intent classification and routing",
            "Six specialist agents each own a health domain:",
            ">> wellness · fitness · dietitian · mental_health · maternal_health · women_health",
            "Each specialist gets only the tools it needs — no bleed between domains",
            "Supervisor receives specialist results and synthesizes the final answer",
            "Benefits: modularity, testability, domain-safe tool scoping, explainability",
        ],
        footer="Defined in react_executor.py — DOMAIN_TOOLSETS and DOMAIN_INSTRUCTIONS",
    )

    # ── SLIDE 7: Section — Code Walkthrough ───────────────────────────────────
    add_section_divider(prs, "SECTION 02", "Code Walkthrough", "LangGraph · State · Nodes · Agent tools")

    # ── SLIDE 8: LangGraph State ───────────────────────────────────────────────
    add_code_slide(
        prs,
        "LangGraph State — HealthCoachState",
        [
            "# src/healthcoach/graph/state.py",
            "",
            "class HealthCoachState(TypedDict):",
            "    conversation_id : str",
            "    user_id         : str",
            "    user_message    : str          # the question from the user",
            "    user_profile    : dict         # age, gender, goal, zip, diet pref",
            "    context         : dict         # loaded from Cosmos DB / stub",
            "    collected_facts : dict         # domain facts gathered by specialists",
            "    final_response  : str          # synthesized answer returned to UI",
        ],
        note="Every node reads from and writes to this shared typed dict — no globals, no side effects",
    )

    # ── SLIDE 9: Graph Workflow ────────────────────────────────────────────────
    add_code_slide(
        prs,
        "LangGraph Workflow — build_graph()",
        [
            "# src/healthcoach/graph/workflow.py",
            "",
            "from langgraph.graph import END, StateGraph",
            "from healthcoach.graph.nodes import load_context, supervisor_agent",
            "",
            "def build_graph():",
            "    graph = StateGraph(HealthCoachState)",
            "    graph.add_node('load_context', load_context)   # enrich state",
            "    graph.add_node('supervisor',   supervisor_agent) # route + orchestrate",
            "    graph.set_entry_point('load_context')",
            "    graph.add_edge('load_context', 'supervisor')",
            "    graph.add_edge('supervisor', END)",
            "    return graph.compile()",
        ],
        note="Deterministic execution: load_context → supervisor → END. Extendable to conditional branches.",
    )

    # ── SLIDE 10: Domain Toolsets ──────────────────────────────────────────────
    add_code_slide(
        prs,
        "Specialist Tool Scoping — DOMAIN_TOOLSETS",
        [
            "# src/healthcoach/agents/react_executor.py",
            "",
            "DOMAIN_TOOLSETS = {",
            "    'wellness'       : ('sleep_summary', 'wellness_docs', 'user_context'),",
            "    'fitness'        : ('workout_summary', 'sleep_summary', 'user_context'),",
            "    'dietitian'      : ('nutrition_docs', 'user_context'),",
            "    'mental_health'  : ('wellness_docs', 'user_context'),",
            "    'maternal_health': ('nutrition_docs', 'wellness_docs', 'user_context'),",
            "    'women_health'   : ('wellness_docs', 'nutrition_docs', 'user_context'),",
            "}",
        ],
        note="Each specialist only sees its scoped tools — prevents hallucination bleed between domains",
    )

    # ── SLIDE 11: ReAct Agent creation ────────────────────────────────────────
    add_code_slide(
        prs,
        "Creating a ReAct Specialist Agent",
        [
            "# src/healthcoach/agents/react_executor.py",
            "",
            "from langgraph.prebuilt import create_react_agent",
            "",
            "def run_domain_react(domain: str, state: HealthCoachState):",
            "    llm          = get_chat_model()           # Azure OpenAI gpt-4o",
            "    tools        = _build_tools(state)        # all available tools",
            "    scoped_tools = [t for t in tools          # filter to domain only",
            "                    if t.name in DOMAIN_TOOLSETS[domain]]",
            "    agent = create_react_agent(",
            "        model=llm, tools=scoped_tools, prompt=system_prompt",
            "    )",
            "    result = agent.invoke({'messages': [HumanMessage(content=user_msg)]})",
        ],
        note="LangGraph prebuilt create_react_agent handles the Reason → Act → Observe loop automatically",
    )

    # ── SLIDE 12: Section — Streamlit UI ──────────────────────────────────────
    add_section_divider(prs, "SECTION 03", "Streamlit UI & User Profile", "app.py — profile capture · state · chat flow")

    # ── SLIDE 13: UI walkthrough ───────────────────────────────────────────────
    add_bullet_slide(
        prs,
        "Streamlit UI — app.py",
        [
            "Sidebar captures: User ID, Gender, Age, Zip/Pincode, Goal, Dietary Preference",
            "Chat input sends query to LangGraph graph via graph.invoke(initial_state)",
            "Initial state created with create_initial_state() — merges UI values into typed dict",
            "Agent response displayed in Streamlit chat_message block",
            "'Show agent facts' toggle reveals collected facts + evaluation scores",
            "Latency timer shows end-to-end response time in milliseconds",
        ],
        footer="app.py — uses @st.cache_resource to build graph once per session",
    )

    # ── SLIDE 14: Section — Azure Integrations ────────────────────────────────
    add_section_divider(prs, "SECTION 04", "Azure Integrations", "OpenAI · AI Search · Cosmos DB · Storage")

    # ── SLIDE 15: Azure services ───────────────────────────────────────────────
    add_bullet_slide(
        prs,
        "Azure Service Integrations",
        [
            "Azure OpenAI (gpt-4o):  powers supervisor and all 6 specialist ReAct agents",
            "Azure AI Search:  semantic retrieval of wellness and nutrition document indexes",
            ">> wellness-index and nutrition-index — grounded evidence for specialist tools",
            "Azure Cosmos DB:  conversation state, user profiles, workout logs",
            "Azure Storage (File Share):  raw wellness document hosting before indexing",
            "Application Insights:  telemetry and connection string in .env.example",
            "All credentials via .env — .gitignore excludes .env from commits",
        ],
        footer="Graceful fallback: app runs in demo mode when Azure credentials are not configured",
    )

    # ── SLIDE 16: Section — Evaluation & Tracing ──────────────────────────────
    add_section_divider(prs, "SECTION 05", "Evaluation & LangSmith Tracing", "criteria.py · langsmith_feedback.py · batch eval")

    # ── SLIDE 17: Evaluation scoring ──────────────────────────────────────────
    add_bullet_slide(
        prs,
        "Custom Evaluation Scoring",
        [
            "Evaluation happens after every graph.invoke() — automatic, not manual",
            "criteria.py defines a weighted multi-criterion rubric:",
            ">> relevance · completeness · safety · actionability · personalization",
            "Each criterion is scored 0–1, combined into an overall weighted score",
            "evaluate_agentic_run() returns scores, reasons, and a pass/fail flag",
            "build_langsmith_feedback_payload() maps scores to LangSmith feedback keys",
            "submit_feedback() attaches scores directly to a run_id in LangSmith UI",
        ],
        footer="scripts/run_batch_eval.py — runs 20 fixed cases and writes batch_eval_report.json",
    )

    # ── SLIDE 18: Section — Doctor Lookup ─────────────────────────────────────
    add_section_divider(prs, "SECTION 06", "Provider Lookup", "US NPI Registry + India Overpass/Nominatim APIs")

    # ── SLIDE 19: Doctor lookup ────────────────────────────────────────────────
    add_two_col_slide(
        prs,
        "Nearby Doctor Lookup — US and India",
        "US: NPI Registry API",
        [
            "Input: US zip code (5-digit)",
            "API: CMS NPI Registry (public, no key required)",
            "Returns licensed providers near the zip",
            "Fields: name, specialty, address, NPI number",
            "Used when query suggests a clinical concern",
        ],
        "India: Nominatim + Overpass API",
        [
            "Input: India pincode (6-digit)",
            "Step 1: Nominatim geocodes pincode → lat/lon",
            "Step 2: Overpass QL queries healthcare nodes",
            "Radius: 15 km from the pincode centroid",
            "Returns: clinic/hospital names and locations",
        ],
    )

    # ── SLIDE 20: Section — Live Demo ─────────────────────────────────────────
    add_section_divider(prs, "SECTION 07", "Live Demo", "Query: 'I am a pre-diabetic patient'")

    # ── SLIDE 21: Demo — User request ─────────────────────────────────────────
    add_chat_demo_slide(
        prs,
        "Demo — User Request",
        "i am a pre diabatic patient",
        "Dietitian + Wellness + Fitness specialists",
        [
            "Routing in progress…",
            "  ▸ Supervisor reading intent → routes to dietitian, wellness, fitness",
            "  ▸ Dietitian agent calling nutrition_docs and user_context tools",
            "  ▸ Wellness agent calling wellness_docs and sleep_summary tools",
            "  ▸ Fitness agent calling workout_summary and sleep_summary tools",
            "  ▸ Supervisor synthesizing collected facts into final response",
        ],
    )

    # ── SLIDE 22: Demo — Agent response ───────────────────────────────────────
    add_chat_demo_slide(
        prs,
        "Demo — Agent Response",
        "i am a pre diabatic patient",
        "Dietitian + Wellness + Fitness specialists",
        [
            "Diet:  Low glycaemic foods — whole grains, legumes, non-starchy vegetables.",
            "       Limit refined carbs, sugar-sweetened drinks, and processed snacks.",
            "Exercise:  30 min moderate activity (walking/cycling) 5 days/week.",
            "           Post-meal walks can help reduce blood sugar spikes.",
            "Wellness:  Monitor fasting glucose regularly. Prioritise 7–9 hours sleep.",
            "Safety note:  This is informational only. Please consult a qualified",
            "              healthcare provider or certified diabetes educator.",
        ],
    )

    # ── SLIDE 23: Section — GitHub & Setup ────────────────────────────────────
    add_section_divider(prs, "SECTION 08", "Get the Code & Run Locally", "github.com/pramodklal/AgenticAI")

    # ── SLIDE 24: Setup commands ───────────────────────────────────────────────
    add_code_slide(
        prs,
        "Quickstart — Clone, Configure, Run",
        [
            "# 1. Clone the repository",
            "git clone https://github.com/pramodklal/AgenticAI.git",
            "cd AgenticAI",
            "",
            "# 2. Create virtual environment and install dependencies",
            "conda create -p envaiagent python=3.11 -y",
            "conda activate ./envaiagent",
            "pip install -r requirements.txt",
            "",
            "# 3. Configure environment",
            "copy .env.example .env    # then fill in your Azure credentials",
            "",
            "# 4. Run the app",
            "streamlit run app.py",
        ],
        note="App runs in demo mode with stub data when Azure keys are not configured — safe to explore first",
    )

    # ── SLIDE 25: What to try ──────────────────────────────────────────────────
    add_bullet_slide(
        prs,
        "Things to Try After Cloning",
        [
            "Ask a health question and enable 'Show agent facts' to see specialist routing",
            "Change profile values (age, goal, dietary preference) and compare responses",
            "Enter a US zip code or India pincode — trigger the doctor lookup routing",
            "Run batch evaluation:  python scripts/run_batch_eval.py",
            "Open LangSmith — view run traces, tool calls, and latency breakdowns",
            "Add a new specialist domain in DOMAIN_TOOLSETS and wire a new tool",
            "Star the repo and open issues for questions or feature requests",
        ],
        footer="github.com/pramodklal/AgenticAI  |  Issues and PRs welcome",
    )

    # ── SLIDE 26: Closing ──────────────────────────────────────────────────────
    add_closing_slide(prs)

    return prs


# ── entry point ────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    prs = build_deck()
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Saved: {OUT}")
