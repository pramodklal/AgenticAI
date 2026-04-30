from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from pptx.util import Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "AI_AgenticAI_MyHealthCoach_Fundamentals.pptx"


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_bullet_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()

    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(20)


def add_quote_slide(prs: Presentation, title: str, quote_lines: list[str], footer_note: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title

    quote_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.7),
        Inches(1.4),
        Inches(12.0),
        Inches(4.7),
    )
    quote_box.fill.solid()
    quote_box.fill.fore_color.rgb = RGBColor(245, 248, 255)
    quote_box.line.color.rgb = RGBColor(12, 61, 130)

    tf = quote_box.text_frame
    tf.clear()
    for i, line in enumerate(quote_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(18)

    note = slide.shapes.add_textbox(Inches(0.8), Inches(6.3), Inches(11.8), Inches(0.8))
    ntf = note.text_frame
    ntf.clear()
    np = ntf.paragraphs[0]
    np.text = footer_note
    np.font.size = Pt(14)
    np.font.color.rgb = RGBColor(80, 80, 80)


def add_chat_slide(prs: Presentation, title: str, speaker: str, lines: list[str], is_user: bool) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title

    if is_user:
        fill_rgb = RGBColor(12, 61, 130)
        text_rgb = RGBColor(255, 255, 255)
        speaker_text = f"User: {speaker}"
    else:
        fill_rgb = RGBColor(235, 242, 255)
        text_rgb = RGBColor(34, 34, 34)
        speaker_text = f"Assistant: {speaker}"

    bubble = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.8),
        Inches(1.4),
        Inches(11.8),
        Inches(4.8),
    )
    bubble.fill.solid()
    bubble.fill.fore_color.rgb = fill_rgb
    bubble.line.color.rgb = RGBColor(12, 61, 130)

    tf = bubble.text_frame
    tf.clear()
    p0 = tf.paragraphs[0]
    p0.text = speaker_text
    p0.font.bold = True
    p0.font.size = Pt(20)
    p0.font.color.rgb = text_rgb

    for line in lines:
        p = tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(18)
        p.font.color.rgb = text_rgb


def build_deck() -> Presentation:
    prs = Presentation()

    add_title_slide(
        prs,
        "AI Agent and Agentic AI Fundamentals",
        "Professional Classroom Deck: Concepts, Frameworks, Patterns, and Real Implementation",
    )

    slides = [
        (
            "Learning Objectives",
            [
                "Define AI Agent and Agentic AI with clear architectural boundaries",
                "Understand agent design choices: single-agent, multi-agent, and tool-centric workflows",
                "Learn practical framework roles: LangChain, LangGraph, LangSmith",
                "Evaluate reasoning patterns: ReAct, Reflection, Plan-and-Execute",
                "Analyze a production-style case study using MyHealthCoach",
            ],
        ),
        (
            "Definition: AI Agent",
            [
                "An AI agent is an autonomous software unit that interprets goals and executes actions",
                "Canonical loop: Observe -> Reason -> Decide -> Act -> Evaluate",
                "Actions may include tool calls, API operations, retrieval, and state updates",
                "An agent is measured by task completion quality, not only text quality",
            ],
        ),
        (
            "Definition: Agentic AI",
            [
                "Agentic AI is a system pattern where planning, execution, and adaptation are first-class",
                "It coordinates reasoning, memory/state, tools, policies, and feedback loops",
                "Can be realized with one orchestrator agent or a specialist multi-agent topology",
                "Primary objective is reliable outcome delivery under real-world constraints",
            ],
        ),
        (
            "AI Agent vs Agentic AI",
            [
                "AI Agent = component-level capability (one decision-making unit)",
                "Agentic AI = system-level architecture and behavior model",
                "A chatbot becomes agentic only when it performs intentional multi-step action",
                "Strong agentic systems include observability, risk controls, and evaluation",
            ],
        ),
        (
            "Types of AI Agents",
            [
                "Reactive agents: fast response under low reasoning depth",
                "Deliberative agents: deeper planning and decomposition before execution",
                "Tool-using agents: operational actions via APIs, DBs, and enterprise systems",
                "Multi-agent systems: orchestration plus domain specialists",
                "Human-in-the-loop systems: controlled escalation for safety-critical workflows",
            ],
        ),
        (
            "Foundational Components",
            [
                "Model layer: reasoning and synthesis engine",
                "Tool layer: actionable capabilities beyond text generation",
                "State and memory: durable context across turns and sessions",
                "Orchestration policy: routing, retries, branching, and stop conditions",
                "Guardrails and evaluator: safety, correctness, and reliability checks",
                "Observability: traceability of every decision, tool call, and latency segment",
            ],
        ),
        (
            "LangChain Overview",
            [
                "Application composition framework for prompts, models, messages, and tools",
                "Provides standardized interfaces for rapid prototyping and refactoring",
                "Useful for implementing tool-calling behavior quickly",
                "Best viewed as the component abstraction layer in the stack",
            ],
        ),
        (
            "LangGraph Overview",
            [
                "State-machine and graph orchestration for controlled agent workflows",
                "Supports deterministic branching, loops, checkpoints, and recovery paths",
                "Well-suited for production-grade, auditable execution",
                "Excellent for supervisor-specialist and multi-stage pipelines",
            ],
        ),
        (
            "LangSmith Overview",
            [
                "Observability and evaluation platform for LLM and agent applications",
                "Captures run trees, tool calls, tokens, latency, and error surfaces",
                "Supports experiment tracking, dataset-driven regression, and comparisons",
                "Enables structured feedback attached directly to run identifiers",
            ],
        ),
        (
            "Agent Framework Patterns",
            [
                "ReAct: iterative reasoning with evidence-gathering actions",
                "Reflection: self-critique pass to improve quality and safety",
                "Plan-and-Execute: decomposition first, deterministic execution second",
                "Self-Consistency: multiple reasoning paths with selection strategy",
                "Multi-Agent Collaboration: explicit role-based specialization",
            ],
        ),
        (
            "ReAct Pattern",
            [
                "Ideal when the model must gather evidence before answering",
                "Loop: reason about next step -> call tool -> integrate observation",
                "Improves groundedness, explainability, and action correctness",
                "Tradeoff: more steps can increase latency and inference cost",
            ],
        ),
        (
            "Reflection Pattern",
            [
                "The system generates a draft and then performs a critique pass",
                "Critique checks logic coverage, safety notes, and response clarity",
                "Can use same model with rubric prompt or a separate critic model",
                "Tradeoff: higher confidence output versus additional compute",
            ],
        ),
        (
            "Case Study Chat: Request and Response",
            [
                "Request: i am a pre diabatic patient",
                "System response (structured): Dietary Recommendations + Exercise Recommendations + Wellness Practices",
                "Safety behavior: explicit healthcare-provider consultation note included",
                "Optional next action: offer to find nearby healthcare provider",
                "Why this is agentic: supervisor routes to dietitian, fitness, and wellness specialists before synthesis",
            ],
        ),
        (
            "MyHealthCoach: Solution Overview",
            [
                "Streamlit UI captures profile, query, and runtime debug controls",
                "LangGraph executes deterministic flow: load_context -> supervisor -> final response",
                "Supervisor ReAct chooses specialist tools based on user intent",
                "Specialists call context, summary, and retrieval tools for grounded advice",
                "Doctor lookup supports US zip and India pincode routing",
            ],
        ),
        (
            "MyHealthCoach: Implementation Details",
            [
                "Supervisor plus six specialist domains: wellness, fitness, dietitian, mental, maternal, women",
                "Doctor router strategy: NPI Registry (US) and Nominatim+Overpass (India)",
                "Azure OpenAI powers both supervisor and specialist reasoning",
                "Data integration path includes Cosmos DB, AI Search, and Azure Storage",
                "Safety-first response policy enforces non-diagnostic advisory tone",
            ],
        ),
        (
            "Evaluation and Performance in This Project",
            [
                "Live per-run scoring appears in Streamlit debug panel",
                "Core criteria: routing accuracy, grounding, actionability, safety, lookup behavior, latency",
                "Batch regression script runs fixed prompt suites for stability checks",
                "LangSmith mapping converts evaluation checks to feedback keys",
                "Optional run_id submission enables experiment-level trend analysis",
            ],
        ),
        (
            "Teaching Lab Activities",
            [
                "Lab 1: convert a single-turn bot to tool-using ReAct",
                "Lab 2: add reflection and compare quality versus latency",
                "Lab 3: create rubric-based evaluation prompts for one health domain",
                "Lab 4: inspect LangSmith traces and identify routing inefficiencies",
                "Lab 5: define safety boundaries for healthcare assistant behavior",
            ],
        ),
        (
            "Key Takeaways",
            [
                "AI Agent is the execution unit; Agentic AI is the orchestration behavior",
                "Framework stack alignment: LangChain + LangGraph + LangSmith",
                "Pattern choice (ReAct/Reflection) should match task risk and complexity",
                "Evaluation and observability are essential for production trust",
                "MyHealthCoach demonstrates a practical agentic architecture in healthcare coaching",
            ],
        ),
    ]

    for title, bullets in slides:
        add_bullet_slide(prs, title, bullets)

    add_chat_slide(
        prs,
        "Case Chat (Request)",
        "i am a pre diabatic patient",
        [
            "Context provided via profile: age, goal, diet preference, and zip/pincode.",
            "Expected behavior: route to dietitian, fitness, and wellness specialists.",
        ],
        is_user=True,
    )

    add_chat_slide(
        prs,
        "Case Chat (Response - Verbatim, Shortened)",
        "Comprehensive plan",
        [
            '"Here is a comprehensive plan tailored to your prediabetes, fat-loss, and energy consistency goals."',
            "Dietary Recommendations: low-glycemic carbs, lean protein, healthy fats, hydration.",
            "Exercise Recommendations: brisk walk, moderate cardio, and recovery-aware strength training.",
            "Wellness Practices: quality sleep, stress management, and routine consistency.",
            '"Safety Note: Consult a healthcare provider before making significant dietary changes."',
            '"Would you like assistance finding a healthcare provider?"',
        ],
        is_user=False,
    )

    return prs


if __name__ == "__main__":
    presentation = build_deck()
    OUT.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(str(OUT))
    print(f"Created: {OUT}")
