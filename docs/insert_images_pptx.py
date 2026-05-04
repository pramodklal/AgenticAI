"""
Inserts logo.png and banner.png into AI_AgenticAI_MyHealthCoach_Fundamentals.pptx:

  Slide  1  (Title)                     → logo.png  bottom-left  (branding)
  Slide 15  (MyHealthCoach Overview)    → banner.png full-width below title
  Slide 21  (last / closing demo)       → logo.png  bottom-centre (closing stamp)

Run: python docs/insert_images_pptx.py
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
import copy

ROOT   = Path(__file__).resolve().parent.parent
PPTX   = ROOT / "docs" / "AI_AgenticAI_MyHealthCoach_Fundamentals.pptx"
LOGO   = ROOT / "assets" / "logo.png"
BANNER = ROOT / "assets" / "banner.png"

prs = Presentation(str(PPTX))

SW = prs.slide_width   # 9144000 EMU = 10 in
SH = prs.slide_height  # 6858000 EMU = 7.5 in


# ── helper ─────────────────────────────────────────────────────────────────────
def add_image(slide, img_path, left, top, width, height=None):
    if height:
        slide.shapes.add_picture(str(img_path), left, top, width, height)
    else:
        slide.shapes.add_picture(str(img_path), left, top, width)


# ── Slide 1 : Title — logo bottom-left ────────────────────────────────────────
slide1 = prs.slides[0]
LOGO_W = Inches(2.4)      # ~240 px wide scaled to slide
LOGO_H = Inches(0.66)
MARGIN = Inches(0.25)
add_image(
    slide1, LOGO,
    left   = MARGIN,
    top    = SH - LOGO_H - MARGIN,
    width  = LOGO_W,
    height = LOGO_H,
)
print("✔  Slide 1  — logo.png added (bottom-left)")


# ── Slide 15 : Solution Overview — banner full-width below title ───────────────
slide15 = prs.slides[14]

# Find title placeholder bottom edge so we position banner just below it
title_bottom = Inches(1.55)   # typical title ends ~1.5 in from top
for sh in slide15.shapes:
    if sh.has_text_frame and sh.top < Inches(1.8):
        candidate = sh.top + sh.height
        if candidate > title_bottom:
            title_bottom = candidate

GAP       = Inches(0.12)
BANNER_W  = SW - Inches(0.5)           # almost full width, small margin each side
BANNER_H  = Inches(1.55)               # maintain ~900:200 aspect roughly
BANNER_L  = Inches(0.25)
BANNER_T  = title_bottom + GAP

add_image(
    slide15, BANNER,
    left   = BANNER_L,
    top    = BANNER_T,
    width  = BANNER_W,
    height = BANNER_H,
)
print("✔  Slide 15 — banner.png added (below title)")


# ── Slide 21 : Last slide — logo bottom-centre ────────────────────────────────
slide21 = prs.slides[20]
LOGO_W2 = Inches(2.6)
LOGO_H2 = Inches(0.72)
MARGIN2 = Inches(0.3)
add_image(
    slide21, LOGO,
    left   = (SW - LOGO_W2) // 2,
    top    = SH - LOGO_H2 - MARGIN2,
    width  = LOGO_W2,
    height = LOGO_H2,
)
print("✔  Slide 21 — logo.png added (bottom-centre)")


# ── save ──────────────────────────────────────────────────────────────────────
prs.save(str(PPTX))
print(f"\nSaved: {PPTX}")
