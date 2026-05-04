from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

PPT_PATH = "docs/MyHealthCoach_YouTube_Video.pptx"

prs = Presentation(PPT_PATH)

# 1) Fix slide 6 text: six -> seven and include medicines.
slide6 = prs.slides[5]
for shape in slide6.shapes:
    if not shape.has_text_frame:
        continue
    tf = shape.text_frame
    for p in tf.paragraphs:
        text = p.text
        if "Six specialist agents each own a health domain" in text:
            p.text = "▸  Seven specialist agents each own a health domain:"
        if "wellness · fitness · dietitian · mental_health · maternal_health · women_health" in text:
            p.text = "    wellness · fitness · dietitian · medicines · mental_health · maternal_health · women_health"

# 2) Create a new slide with all 8 agents and short descriptions.
blank_layout = prs.slide_layouts[6]
new_slide = prs.slides.add_slide(blank_layout)

# Colors aligned with existing deck style.
DARK = RGBColor(0x0D, 0x1B, 0x2A)
ORANGE = RGBColor(0xF5, 0x6E, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0xE5, 0xE7, 0xEB)
CARD = RGBColor(0x1A, 0x3A, 0x5C)

# Background + accent bar.
bg = new_slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
bg.fill.solid()
bg.fill.fore_color.rgb = DARK
bg.line.fill.background()

bar = new_slide.shapes.add_shape(1, 0, 0, prs.slide_width, Pt(6))
bar.fill.solid()
bar.fill.fore_color.rgb = ORANGE
bar.line.fill.background()

# Title
box_title = new_slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12), Inches(0.6))
tf_title = box_title.text_frame
p = tf_title.paragraphs[0]
r = p.add_run()
r.text = "All 8 Agents in MyHealthCoach"
r.font.bold = True
r.font.size = Pt(30)
r.font.color.rgb = ORANGE

# Subtitle
box_sub = new_slide.shapes.add_textbox(Inches(0.5), Inches(0.85), Inches(12), Inches(0.35))
tf_sub = box_sub.text_frame
p = tf_sub.paragraphs[0]
r = p.add_run()
r.text = "Each agent has a focused role and the supervisor coordinates the final response."
r.font.size = Pt(13)
r.font.color.rgb = MUTED

agents = [
    ("Supervisor Agent", "Routes intent to the right specialists and combines outputs.", "Ensures grounded, safe, and structured final response."),
    ("Wellness Agent", "Handles sleep, recovery, stress, and fatigue guidance.", "Uses user context and wellness knowledge for daily habits."),
    ("Fitness Agent", "Handles workouts, training load, and movement safety.", "Suggests exercise adjustments based on recovery readiness."),
    ("Dietitian Agent", "Handles nutrition, hydration, and meal-planning advice.", "Provides practical, goal-aligned dietary recommendations."),
    ("Medicines Agent", "Handles symptom-aware OTC guidance and medicine safety.", "Adds red-flag escalation when serious symptom patterns appear."),
    ("Mental Health Agent", "Handles anxiety, mood support, and stress de-escalation.", "Suggests supportive routines and when to seek licensed help."),
    ("Maternal Health Agent", "Handles prenatal and postpartum-safe guidance.", "Emphasizes low-risk routines with explicit escalation notes."),
    ("Women Health Agent", "Handles cycle-related wellness, PCOS, and menopause topics.", "Adapts activity and nutrition to symptom patterns."),
]

# 2 columns x 4 rows cards
left_x = Inches(0.5)
right_x = Inches(6.55)
start_y = Inches(1.35)
card_w = Inches(5.65)
card_h = Inches(1.25)
gap_y = Inches(0.15)

for i, (name, line1, line2) in enumerate(agents):
    col = 0 if i < 4 else 1
    row = i if i < 4 else i - 4
    x = left_x if col == 0 else right_x
    y = start_y + row * (card_h + gap_y)

    card = new_slide.shapes.add_shape(1, x, y, card_w, card_h)
    card.fill.solid()
    card.fill.fore_color.rgb = CARD
    card.line.color.rgb = ORANGE

    t = new_slide.shapes.add_textbox(x + Inches(0.18), y + Inches(0.10), card_w - Inches(0.3), card_h - Inches(0.18))
    tf = t.text_frame
    tf.clear()

    p0 = tf.paragraphs[0]
    r0 = p0.add_run()
    r0.text = name
    r0.font.bold = True
    r0.font.size = Pt(13)
    r0.font.color.rgb = ORANGE

    p1 = tf.add_paragraph()
    r1 = p1.add_run()
    r1.text = line1
    r1.font.size = Pt(11)
    r1.font.color.rgb = WHITE

    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = line2
    r2.font.size = Pt(11)
    r2.font.color.rgb = WHITE

# 3) Move new slide to position after slide 6 (index 6, i.e., new slide 7)
slides_xml = prs.slides._sldIdLst
items = list(slides_xml)
new_elem = items[-1]
slides_xml.remove(new_elem)
slides_xml.insert(6, new_elem)

prs.save(PPT_PATH)
print("Updated slide 6 and inserted new slide 7.")
