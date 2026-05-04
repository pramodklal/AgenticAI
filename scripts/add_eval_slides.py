"""Add 4 evaluation-score slides to MyHealthCoach_YouTube_Video.pptx."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

PPTX_PATH = "docs/MyHealthCoach_YouTube_Video.pptx"

prs = Presentation(PPTX_PATH)
W = prs.slide_width
H = prs.slide_height

DARK_BLUE  = RGBColor(0x0D, 0x1B, 0x2A)
ORANGE     = RGBColor(0xF5, 0x6E, 0x00)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xF2, 0xF4, 0xF7)
GREEN      = RGBColor(0x00, 0x96, 0x47)
AMBER      = RGBColor(0xFF, 0xBF, 0x00)
MID_BLUE   = RGBColor(0x1A, 0x3A, 0x5C)
CODE_BG    = RGBColor(0x0A, 0x14, 0x1E)
CYAN       = RGBColor(0x56, 0xB6, 0xC2)
DARK_GREEN = RGBColor(0x00, 0x5C, 0x2E)
DARK_AMBER = RGBColor(0x7A, 0x50, 0x00)


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def bg_rect(slide, color):
    s = slide.shapes.add_shape(1, 0, 0, W, H)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def accent_bar(slide, color=None):
    color = color or ORANGE
    s = slide.shapes.add_shape(1, 0, 0, W, Pt(6))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def box(slide, color, left, top, width, height):
    s = slide.shapes.add_shape(1, left, top, width, height)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def tb(slide, text, left, top, width, height,
       font_size=18, bold=False, color=WHITE,
       align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(left, top, width, height)
    shape.text_frame.word_wrap = True
    p = shape.text_frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def move_slide(prs, old_index, new_index):
    xml_slides = prs.slides._sldIdLst
    items = list(xml_slides)
    elem = items[old_index]
    xml_slides.remove(elem)
    xml_slides.insert(new_index, elem)


# ════════════════════════════════════════════════════════
# SLIDE A — The 6 Evaluation Checks
# ════════════════════════════════════════════════════════
sA = blank_slide(prs)
bg_rect(sA, DARK_BLUE)
accent_bar(sA)

tb(sA, "Evaluation — The 6 Checks",
   Inches(0.5), Inches(0.12), Inches(12), Inches(0.7),
   font_size=30, bold=True, color=ORANGE)

tb(sA, "Each check has an independent pass threshold and weight in the weighted average.",
   Inches(0.5), Inches(0.82), Inches(11.5), Inches(0.4),
   font_size=15, color=LIGHT_GREY)

CHECKS = [
    ("domain_routing_accuracy",  "Did the supervisor call the RIGHT specialist tools?",          "22%", ">= 0.70"),
    ("tool_usage_grounding",     "Is the answer grounded in tool output — not LLM memory?",      "15%", ">= 0.60"),
    ("response_actionability",   "Is the advice specific, practical, and actionable?",            "14%", ">= 0.60"),
    ("safety_and_escalation",    "Does the response include safety warnings and escalation cues?","20%", ">= 0.75"),
    ("doctor_lookup_behavior",   "Did find_nearby_doctors run and return structured results?",    "14%", ">= 0.70"),
    ("latency_budget",           "Did the full run finish within the latency budget?",            "15%", ">= 0.70"),
]

col_x = [Inches(0.40), Inches(3.80), Inches(10.55), Inches(11.50)]
col_w = [Inches(3.30), Inches(6.65), Inches(0.90),  Inches(0.90)]
row_h = Inches(0.52)
start_y = Inches(1.35)

# header
for text, cx, cw in zip(["CHECK KEY", "WHAT IT MEASURES", "WEIGHT", "THRESHOLD"],
                         col_x, col_w):
    tb(sA, text, cx, start_y, cw, row_h, font_size=12, bold=True, color=ORANGE)

for ri, (key, desc, wt, thr) in enumerate(CHECKS):
    y = start_y + row_h + Pt(4) + ri * (row_h + Pt(2))
    box(sA, MID_BLUE if ri % 2 == 0 else DARK_BLUE,
        Inches(0.35), y - Pt(3), W - Inches(0.70), row_h)
    for text, cx, cw in zip([key, desc, wt, thr], col_x, col_w):
        clr = AMBER if text == wt else (GREEN if text == thr else WHITE)
        tb(sA, text, cx, y, cw, row_h, font_size=13,
           bold=(text == key), color=clr)

tb(sA, "Overall Score  =  Sum(check_score x weight)   |   Passed when overall >= 0.75",
   Inches(0.5), Inches(6.1), Inches(11.5), Inches(0.5),
   font_size=14, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════
# SLIDE B — Score Formula Walkthrough
# ════════════════════════════════════════════════════════
sB = blank_slide(prs)
bg_rect(sB, DARK_BLUE)
accent_bar(sB)

tb(sB, "How the Weighted Score Is Calculated",
   Inches(0.5), Inches(0.12), Inches(12), Inches(0.7),
   font_size=30, bold=True, color=ORANGE)

# Left formula panel
box(sB, MID_BLUE, Inches(0.40), Inches(0.95), Inches(5.60), Inches(4.80))

tb(sB, "Overall Score  =  Sum(score_i  x  weight_i)",
   Inches(0.55), Inches(1.00), Inches(5.30), Inches(0.55),
   font_size=15, bold=True, color=ORANGE)

FORMULA_ROWS = [
    ("domain_routing",    "1.0  x  0.22  =  0.220", True),
    ("tool_grounding",    "1.0  x  0.15  =  0.150", True),
    ("actionability",     "1.0  x  0.14  =  0.140", True),
    ("safety_escalation", "1.0  x  0.20  =  0.200", True),
    ("doctor_lookup",     "1.0  x  0.14  =  0.140", True),
    ("latency_budget",    "0.85 x  0.15  =  0.128", False),
]
for fi, (name, calc, perfect) in enumerate(FORMULA_ROWS):
    fy = Inches(1.65) + fi * Inches(0.57)
    tb(sB, name, Inches(0.60), fy, Inches(2.60), Inches(0.50),
       font_size=13, color=LIGHT_GREY)
    tb(sB, calc, Inches(3.25), fy, Inches(2.50), Inches(0.50),
       font_size=13, bold=True, color=GREEN if perfect else AMBER)

box(sB, ORANGE, Inches(0.55), Inches(5.10), Inches(4.90), Pt(2))
tb(sB, "TOTAL  =  0.978     Passed (>= 0.75)",
   Inches(0.55), Inches(5.20), Inches(5.20), Inches(0.55),
   font_size=16, bold=True, color=GREEN)

# Right latency panel
tb(sB, "Why only 0.85 on latency?",
   Inches(6.40), Inches(0.95), Inches(5.50), Inches(0.55),
   font_size=18, bold=True, color=ORANGE)

box(sB, MID_BLUE, Inches(6.40), Inches(1.60), Inches(5.50), Inches(3.30))

LAT_LINES = [
    ("Target:    <= 8,000 ms",                              WHITE),
    ("Actual:    11,491 ms  (two sequential agent calls)",  AMBER),
    ("",                                                    WHITE),
    ("score = max(0,  1 - (elapsed - target) / 10000)",    CYAN),
    ("      = max(0,  1 - (11491 - 8000) / 10000)",        CYAN),
    ("      = max(0,  1 - 0.3491)",                        CYAN),
    ("      = 0.85",                                        GREEN),
    ("",                                                    WHITE),
    ("Still PASSES — threshold is only 0.70",               GREEN),
]
for li, (line, clr) in enumerate(LAT_LINES):
    tb(sB, line, Inches(6.60), Inches(1.70) + li * Inches(0.33),
       Inches(5.20), Inches(0.32), font_size=12,
       bold=("PASSES" in line), color=clr)

tb(sB, ("Why slow?  Two sequential calls:\n"
        "  1.  consult_medicines_advisor  (Azure OpenAI + AI Search)\n"
        "  2.  find_nearby_doctors  (NPI Registry HTTP)"),
   Inches(6.40), Inches(5.05), Inches(5.50), Inches(0.95),
   font_size=13, color=LIGHT_GREY)

# ════════════════════════════════════════════════════════
# SLIDE C — Live Run Result (0.978)
# ════════════════════════════════════════════════════════
sC = blank_slide(prs)
bg_rect(sC, DARK_BLUE)
accent_bar(sC)

tb(sC, "Live Run — Eye Flu Query   |   Score: 0.978   PASS",
   Inches(0.5), Inches(0.12), Inches(12), Inches(0.7),
   font_size=28, bold=True, color=ORANGE)

tb(sC, ('Query: "I am having eye flu symptoms"'
        "   |   zip_code: 08816"
        "   |   Latency: 11,491 ms"),
   Inches(0.5), Inches(0.82), Inches(11.5), Inches(0.38),
   font_size=14, color=LIGHT_GREY)

# Left — collected_facts snapshot
box(sC, MID_BLUE, Inches(0.35), Inches(1.28), Inches(6.10), Inches(5.30))

tb(sC, "collected_facts  (medicines domain)",
   Inches(0.50), Inches(1.33), Inches(5.80), Inches(0.42),
   font_size=13, bold=True, color=ORANGE)

FACTS = [
    ("recommendation:", "Use artificial tears / lubricating eye drops.", True),
    ("",                "Cold compresses to reduce swelling.",            False),
    ("",                "Wash hands. Avoid touching eyes.",              False),
    ("rationale:",      "Artificial tears alleviate dryness.",           True),
    ("",                "Cold compresses reduce discomfort.",            False),
    ("",                "Hygiene prevents further spread.",              False),
    ("safety_note:",    "Avoid steroid/antibiotic drops without doctor.", True),
    ("",                "If symptoms persist >3 days see a doctor.",     False),
    ("confidence:",     "medium",                                        True),
]
for fi, (label, val, is_label_row) in enumerate(FACTS):
    fy = Inches(1.82) + fi * Inches(0.48)
    if label:
        tb(sC, label, Inches(0.50), fy, Inches(1.55), Inches(0.45),
           font_size=11, bold=True, color=AMBER)
    tb(sC, val, Inches(1.90), fy, Inches(4.40), Inches(0.45),
       font_size=11, color=WHITE)

# Right — score summary
box(sC, MID_BLUE, Inches(6.60), Inches(1.28), Inches(5.20), Inches(5.30))

tb(sC, "Evaluation Result",
   Inches(6.75), Inches(1.33), Inches(4.80), Inches(0.42),
   font_size=13, bold=True, color=ORANGE)

SCORE_ROWS = [
    ("domain_routing_accuracy",  "1.0",  True),
    ("tool_usage_grounding",     "1.0",  True),
    ("response_actionability",   "1.0",  True),
    ("safety_and_escalation",    "1.0",  True),
    ("doctor_lookup_behavior",   "1.0",  True),
    ("latency_budget",           "0.85", False),
]
for si, (key, score, perfect) in enumerate(SCORE_ROWS):
    sy = Inches(1.82) + si * Inches(0.62)
    box(sC, DARK_GREEN if perfect else DARK_AMBER,
        Inches(6.65), sy - Pt(2), Inches(5.10), Inches(0.55))
    tb(sC, key, Inches(6.75), sy, Inches(3.60), Inches(0.50),
       font_size=12, color=WHITE)
    tb(sC, score, Inches(10.45), sy, Inches(0.90), Inches(0.50),
       font_size=13, bold=True,
       color=GREEN if perfect else AMBER, align=PP_ALIGN.RIGHT)

box(sC, ORANGE, Inches(6.65), Inches(5.65), Inches(5.00), Pt(2))
tb(sC, "OVERALL  =  0.978     PASSED",
   Inches(6.75), Inches(5.75), Inches(4.80), Inches(0.55),
   font_size=17, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════
# SLIDE D — Code Implementation
# ════════════════════════════════════════════════════════
sD = blank_slide(prs)
bg_rect(sD, DARK_BLUE)
accent_bar(sD)

tb(sD, "How It Is Implemented — criteria.py & react_executor.py",
   Inches(0.5), Inches(0.12), Inches(12), Inches(0.7),
   font_size=26, bold=True, color=ORANGE)

# Left code panel
box(sD, CODE_BG, Inches(0.35), Inches(0.92), Inches(5.85), Inches(5.70))

tb(sD, "src/healthcoach/evaluation/criteria.py",
   Inches(0.50), Inches(0.97), Inches(5.60), Inches(0.38),
   font_size=12, bold=True, color=AMBER)

CODE_LEFT = [
    ("# 6 weighted checks",                                            AMBER),
    ("CHECKS = [",                                                     WHITE),
    ('  EvalCheck("domain_routing_accuracy",  wt=0.22, thr=0.70),',   WHITE),
    ('  EvalCheck("tool_usage_grounding",     wt=0.15, thr=0.60),',   WHITE),
    ('  EvalCheck("response_actionability",   wt=0.14, thr=0.60),',   WHITE),
    ('  EvalCheck("safety_and_escalation",    wt=0.20, thr=0.75),',   WHITE),
    ('  EvalCheck("doctor_lookup_behavior",   wt=0.14, thr=0.70),',   WHITE),
    ('  EvalCheck("latency_budget",           wt=0.15, thr=0.70),',   WHITE),
    ("]",                                                              WHITE),
    ("",                                                               WHITE),
    ("# Latency scoring formula",                                      AMBER),
    ("def _score_latency(elapsed_ms):",                                CYAN),
    ("    target = 8_000",                                             WHITE),
    ("    return max(0.0, 1-(elapsed_ms-target)/10_000)",              WHITE),
    ("",                                                               WHITE),
    ("# Fever/symptom paired domain routing",                          AMBER),
    ("if any(h in query for h in _FEVER_LIKE_HINTS):",                WHITE),
    ('    expected |= {"wellness", "medicines"}',                      GREEN),
]
for ci, (line, clr) in enumerate(CODE_LEFT):
    tb(sD, line, Inches(0.50), Inches(1.42) + ci * Inches(0.286),
       Inches(5.60), Inches(0.28), font_size=10, color=clr)

# Right code panel
box(sD, CODE_BG, Inches(6.40), Inches(0.92), Inches(5.50), Inches(5.70))

tb(sD, "src/healthcoach/agents/react_executor.py",
   Inches(6.50), Inches(0.97), Inches(5.30), Inches(0.38),
   font_size=12, bold=True, color=AMBER)

CODE_RIGHT = [
    ("# Supervisor MUST call tools before answering",                  AMBER),
    ("system_prompt = (",                                              WHITE),
    ('  "IMPORTANT: You MUST always call at least"',                   WHITE),
    ('  " one specialist advisor tool before answering."',             WHITE),
    ('  " Never answer from your own knowledge."',                     WHITE),
    ("  ...",                                                          LIGHT_GREY),
    (")",                                                              WHITE),
    ("",                                                               WHITE),
    ("# Preserve full advisor JSON — no truncation",                   AMBER),
    ("is_consult = tool_name.startswith('consult_')",                  WHITE),
    ("observation = (content if is_consult",                          WHITE),
    ("               else content[:2000])",                           WHITE),
    ("",                                                               WHITE),
    ("# Rebuild collected_facts from tool trace",                      AMBER),
    ("for item in tool_trace:",                                        WHITE),
    ("    domain = (item['tool']",                                     WHITE),
    ("              .replace('consult_','')",                         WHITE),
    ("              .replace('_advisor',''))",                        WHITE),
    ("    collected_facts[domain] = json.loads(",                      WHITE),
    ("        item['observation'])",                                   WHITE),
]
for ci, (line, clr) in enumerate(CODE_RIGHT):
    tb(sD, line, Inches(6.50), Inches(1.42) + ci * Inches(0.246),
       Inches(5.30), Inches(0.26), font_size=10, color=clr)

# ════════════════════════════════════════════════════════
# Reorder: move the 4 new slides to positions 19-22
# (right after existing slide 19 = index 18)
# ════════════════════════════════════════════════════════
# New slides are currently at the end: indices 28,29,30,31
# Insert each in turn starting at target position 19
total = len(prs.slides)  # should be 32
for i in range(4):
    move_slide(prs, total - 1, 19 + i)

prs.save(PPTX_PATH)
print(f"Saved. Total slides: {len(prs.slides)}")
